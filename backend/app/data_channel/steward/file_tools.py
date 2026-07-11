"""Safe document creation for a single data-steward conversation.

The model never supplies a filesystem path.  It supplies a display filename and
content; this module renders bytes in memory and hands them to ``workspace``,
which is the only component allowed to resolve a conversation directory.
"""
from __future__ import annotations

import csv
import io
import re
from pathlib import Path
from typing import Any

from app.data_channel.steward import workspace


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf", ".md", ".txt", ".csv"}
MIME_TYPES = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pdf": "application/pdf",
    ".md": "text/markdown; charset=utf-8",
    ".txt": "text/plain; charset=utf-8",
    ".csv": "text/csv; charset=utf-8",
}
_MAX_GENERATED_CHARS = 500_000


def _filename(filename: str) -> tuple[str, str]:
    safe = workspace.safe_filename(filename, "document.md")
    suffix = Path(safe).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        supported = "、".join(sorted(SUPPORTED_EXTENSIONS))
        raise workspace.WorkspaceError(f"不支持创建 {suffix or '无扩展名'} 文件；支持：{supported}")
    return safe, suffix


def _text(content: Any) -> str:
    if content is None:
        return ""
    value = content if isinstance(content, str) else str(content)
    if len(value) > _MAX_GENERATED_CHARS:
        raise workspace.WorkspaceError(f"生成内容超过 {_MAX_GENERATED_CHARS} 字符限制")
    return value


def _paragraphs(content: str) -> list[str]:
    return [part.strip() for part in re.split(r"\n\s*\n", content) if part.strip()] or [""]


def _render_docx(content: str, title: str | None) -> bytes:
    from docx import Document

    document = Document()
    if title:
        document.add_heading(title.strip(), level=0)
    for block in _paragraphs(content):
        lines = block.splitlines()
        first = lines[0].strip()
        heading = re.match(r"^(#{1,3})\s+(.+)$", first)
        if heading:
            document.add_heading(heading.group(2), level=len(heading.group(1)))
            lines = lines[1:]
        for line in lines:
            value = line.strip()
            if value.startswith(("- ", "* ")):
                document.add_paragraph(value[2:], style="List Bullet")
            elif re.match(r"^\d+[.)]\s+", value):
                document.add_paragraph(re.sub(r"^\d+[.)]\s+", "", value), style="List Number")
            elif value:
                document.add_paragraph(value)
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _render_pptx(content: str, title: str | None) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    blocks = _paragraphs(content)
    if title:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title.strip()
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "由数据管家在当前会话中生成"
    for index, block in enumerate(blocks):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        heading = re.sub(r"^#{1,6}\s+", "", lines[0]) if lines else f"第 {index + 1} 页"
        slide.shapes.title.text = heading[:120]
        body = slide.placeholders[1].text_frame
        body.clear()
        values = lines[1:] if len(lines) > 1 else lines
        for item_index, line in enumerate(values or [""]):
            paragraph = body.paragraphs[0] if item_index == 0 else body.add_paragraph()
            paragraph.text = re.sub(r"^(?:[-*]|\d+[.)])\s+", "", line)
            paragraph.font.size = Pt(20)
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _table_rows(content: str, rows: list[Any] | None) -> list[list[Any]]:
    if rows:
        if all(isinstance(row, dict) for row in rows):
            headers: list[str] = []
            for row in rows:
                for key in row:
                    if str(key) not in headers:
                        headers.append(str(key))
            return [headers] + [[row.get(key, "") for key in headers] for row in rows]
        if all(isinstance(row, (list, tuple)) for row in rows):
            return [list(row) for row in rows]
        raise workspace.WorkspaceError("rows 必须是对象数组或二维数组")
    parsed: list[list[str]] = []
    for line in content.splitlines():
        if not line.strip():
            continue
        delimiter = "\t" if "\t" in line else ("," if "," in line else None)
        parsed.append([cell.strip() for cell in line.split(delimiter)] if delimiter else [line.strip()])
    return parsed or [[""]]


def _render_xlsx(content: str, rows: list[Any] | None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    table = _table_rows(content, rows)
    for row in table:
        sheet.append(row)
    if table:
        for cell in sheet[1]:
            cell.font = Font(bold=True)
    for column in sheet.columns:
        width = min(60, max(10, max(len(str(cell.value or "")) for cell in column) + 2))
        sheet.column_dimensions[column[0].column_letter].width = width
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _render_pdf(content: str, title: str | None) -> bytes:
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    from xml.sax.saxutils import escape

    font_name = "STSong-Light"
    pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    output = io.BytesIO()
    document = SimpleDocTemplate(output, pagesize=A4, leftMargin=22 * mm,
                                 rightMargin=22 * mm, topMargin=20 * mm, bottomMargin=20 * mm)
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle("ChineseBody", parent=styles["BodyText"], fontName=font_name,
                                fontSize=11, leading=18, spaceAfter=8)
    title_style = ParagraphStyle("ChineseTitle", parent=styles["Title"], fontName=font_name,
                                 fontSize=20, leading=28, alignment=TA_CENTER, spaceAfter=14)
    story = []
    if title:
        story.extend([Paragraph(escape(title.strip()), title_style), Spacer(1, 4 * mm)])
    for block in _paragraphs(content):
        story.append(Paragraph(escape(block).replace("\n", "<br/>") or " ", body_style))
    document.build(story)
    return output.getvalue()


def render(filename: str, content: Any, *, title: str | None = None,
           rows: list[Any] | None = None) -> tuple[str, bytes, str]:
    safe, suffix = _filename(filename)
    value = _text(content)
    if suffix == ".docx":
        data = _render_docx(value, title)
    elif suffix == ".pptx":
        data = _render_pptx(value, title)
    elif suffix == ".xlsx":
        data = _render_xlsx(value, rows)
    elif suffix == ".pdf":
        data = _render_pdf(value, title)
    elif suffix == ".csv":
        output = io.StringIO(newline="")
        writer = csv.writer(output)
        writer.writerows(_table_rows(value, rows))
        data = output.getvalue().encode("utf-8-sig")
    else:
        data = value.encode("utf-8")
    return safe, data, MIME_TYPES[suffix]


def create(conversation_id: str, filename: str, content: Any, *, title: str | None = None,
           rows: list[Any] | None = None, source: str = "generated") -> dict:
    safe, data, mime = render(filename, content, title=title, rows=rows)
    return workspace.save_bytes(conversation_id, safe, data, source=source,
                                mime_type=mime, extract=True)


def edit(conversation_id: str, artifact_id: str, content: Any, *, mode: str = "replace",
         output_filename: str | None = None, title: str | None = None,
         rows: list[Any] | None = None) -> dict:
    row, _path = workspace.require_file(conversation_id, artifact_id)
    normalized_mode = (mode or "replace").strip().lower()
    if normalized_mode not in {"replace", "append"}:
        raise workspace.WorkspaceError("mode 只能是 replace 或 append")
    old_text = workspace.extracted_text(conversation_id, artifact_id, _MAX_GENERATED_CHARS)
    if normalized_mode == "append":
        new_text = f"{old_text.rstrip()}\n\n{_text(content).lstrip()}".strip()
    else:
        new_text = _text(content)
    original = Path(row["filename"])
    filename = output_filename or f"{original.stem}_edited{original.suffix}"
    created = create(conversation_id, filename, new_text, title=title, rows=rows, source="edited")
    created["editedFromArtifactId"] = artifact_id
    created["editMode"] = normalized_mode
    return created
